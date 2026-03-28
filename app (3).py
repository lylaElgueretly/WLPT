import sys
import subprocess
from pathlib import Path

# --- Runtime package installation to avoid ModuleNotFoundError ---
def install(package):
    subprocess.check_call([sys.executable, "-m", "pip", "install", package])

try:
    from docx import Document
except ModuleNotFoundError:
    install("python-docx")
    from docx import Document

try:
    import PyPDF2
except ModuleNotFoundError:
    install("PyPDF2")
    import PyPDF2

try:
    from pptx import Presentation
except ModuleNotFoundError:
    install("python-pptx")
    from pptx import Presentation

# Streamlit for UI (optional, remove if not using Streamlit)
try:
    import streamlit as st
except ModuleNotFoundError:
    install("streamlit")
    import streamlit as st

# --- Function to extract text from different file types ---
def extract_text(file_path):
    ext = file_path.suffix.lower()
    
    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    
    elif ext == ".docx":
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    
    elif ext == ".pdf":
        text = ""
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text
    
    elif ext == ".pptx":
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text
    
    else:
        return f"[Unsupported file type: {ext}]"

# --- Main App Logic ---
def process_folder(folder_path):
    folder = Path(folder_path)
    for file in folder.iterdir():
        if file.is_file():
            print(f"Processing: {file.name}")
            text = extract_text(file)
            # Replace this line with your app's actual processing logic
            print(f"Extracted Text (first 500 chars):\n{text[:500]}...\n")  # Preview first 500 chars

# --- Streamlit UI (optional) ---
if 'st' in sys.modules:
    st.title("Document Reader App")
    folder_input = st.text_input("Enter folder path to process:", "C:/Users/Lyla/Downloads/myproject")
    if st.button("Process Files"):
        st.text(f"Processing folder: {folder_input}")
        folder = Path(folder_input)
        for file in folder.iterdir():
            if file.is_file():
                text = extract_text(file)
                st.subheader(file.name)
                st.text_area("Extracted Text (first 500 chars)", text[:500])

# --- If running as a normal script ---
if __name__ == "__main__" and 'streamlit' not in sys.modules:
    folder_path = "C:/Users/Lyla/Downloads/myproject"  # Change this to your folder path
    process_folder(folder_path)
