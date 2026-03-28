# app.py
from pathlib import Path
from docx import Document
import PyPDF2
from pptx import Presentation
import streamlit as st

# --- Function to extract text from different file types ---
def extract_text(file_path):
    ext = file_path.suffix.lower()
    
    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    
    elif ext == ".docx":
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    
    elif ext == ".pdf":
        text = ""
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text
    
    elif ext == ".pptx":
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text
    
    else:
        return f"[Unsupported file type: {ext}]"

# --- Function to process all files in a folder ---
def process_folder(folder_path):
    folder = Path(folder_path)
    results = {}
    for file in folder.iterdir():
        if file.is_file():
            text = extract_text(file)
            results[file.name] = text
    return results

# --- Streamlit UI ---
st.title("Document Reader App")

folder_input = st.text_input(
    "Enter folder path to process:",
    "C:/Users/Lyla/Downloads/myproject"
)

if st.button("Process Files"):
    folder_path = Path(folder_input)
    if folder_path.exists() and folder_path.is_dir():
        results = process_folder(folder_path)
        for filename, text in results.items():
            st.subheader(filename)
            st.text_area("Extracted Text (first 500 chars)", text[:500])
    else:
        st.error("Folder path does not exist or is not a directory.")

# --- Optional: Run as standalone script ---
if __name__ == "__main__":
    folder_path = "C:/Users/Lyla/Downloads/myproject"  # Change this path if needed
    results = process_folder(folder_path)
    for filename, text in results.items():
        print(f"--- {filename} ---")
        print(text[:500])  # Preview first 500 chars
        print("\n")
